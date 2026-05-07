"""
Integrations Router — Google Classroom, Drive, Gmail sync + Moodle
Note: These use standard OAuth scopes. No Google app verification needed
when the app is in "Testing" mode (up to 100 test users).
"""
from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks, File, UploadFile
from fastapi.responses import StreamingResponse, RedirectResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from pydantic import BaseModel
from typing import Optional, List, Dict
import asyncio
import json
import logging
import re
import os
from datetime import datetime, timezone, timedelta
from app.models.database import get_db, User, Course, Material
from app.utils.dependencies import get_current_user
from app.services.google_auth import build_credentials, refresh_credentials
from app.services.drive_sync import DriveSyncService
from app.intelligence.ms_client import get_ms_client
import io
import base64

logger = logging.getLogger(__name__)

class DriveSyncRequest(BaseModel):
    file_ids: List[str]

router = APIRouter()

@router.get("/test-upload")
async def test_upload():
    return {"status": "ok", "message": "Integrations router is reachable"}

async def _extract_text_from_bytes(raw_bytes: bytes, mime_type: str) -> str:
    """Helper to extract text from raw bytes based on mime type for manual uploads."""
    try:
        if mime_type == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            text = "\n".join(p.extract_text() or "" for p in reader.pages)[:15000]
            return text
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            doc = docx.Document(io.BytesIO(raw_bytes))
            text = "\n".join(para.text for para in doc.paragraphs)[:15000]
            return text
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                text_runs.append("\n")
            return "".join(text_runs)[:15000]
        elif mime_type.startswith("image/"):
            try:
                b64 = base64.b64encode(raw_bytes).decode("utf-8")
                from app.services.openai_service import unified_completion
                
                # Build messages for vision
                messages = [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract all text from this image and provide a high-level academic description of what it shows. Return the transcription followed by the description."},
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{b64}"}
                            }
                        ]
                    }
                ]
                
                # Use unified completion (LiteLLM handles vision for Gemini/GPT-4o)
                resp = await unified_completion(messages, tier="pro")
                return resp["text"][:15000]
            except Exception as e:
                logger.error(f"Manual vision OCR error: {e}")
                return "[Image content could not be processed]"
        elif mime_type == "text/plain" or mime_type.startswith("text/"):
            return raw_bytes.decode("utf-8", errors="replace")[:15000]
        return ""
    except Exception as e:
        logger.error(f"Manual text extraction failed: {e}")
        return ""


@router.post("/upload")
async def manual_upload(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Manually upload a document to both local DB and Memsapien."""
    # 1. Read file
    file_bytes = await file.read()
    filename = file.filename or "uploaded_file"
    mime_type = file.content_type or "application/octet-stream"
    
    # 2. Extract text
    content = await _extract_text_from_bytes(file_bytes, mime_type)
    preview = content[:500] if content else f"Manually uploaded file: {filename}"
    
    # 3. Save to local DB
    new_mat = Material(
        user_id=current_user.id,
        title=filename,
        material_type="file",
        source="manual_upload",
        file_content=file_bytes,
        mime_type=mime_type,
        indexed_in_memsapien=False,
        content_preview=preview,
        full_text=content if content else None,
        source_link=f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/0"
    )
    db.add(new_mat)
    await db.commit()
    await db.refresh(new_mat)
    
    # Update with real ID link
    backend_url = os.getenv("BACKEND_URL", "http://localhost:8002")
    new_mat.source_link = f"{backend_url}/profile/proxy/drive/{new_mat.id}"
    await db.commit()
    
    # 4. Index in Memsapien
    try:
        sm = get_ms_client()
        sm_id = await sm.add_document(
            content=f"User Uploaded Document: {filename}\n\n{content or 'Binary file content'}",
            metadata={
                "source": "manual_upload",
                "type": "user_upload",
                "user_id": current_user.email,
                "material_id": str(new_mat.id),
                "source_link": new_mat.source_link
            },
            user_email=current_user.email,
            title=filename,
            container_tag=f"user_{current_user.id}"
        )
        if sm_id:
            new_mat.indexed_in_memsapien = True
            await db.commit()
    except Exception as e:
        logger.error(f"Failed to index manual upload in SM: {e}")
        
    return {"id": new_mat.id, "title": filename, "indexed": new_mat.indexed_in_memsapien}



def _extract_classroom_attachments(materials_list: list) -> list:
    """Convert a Classroom API materials list into a JSON-serializable attachments array."""
    attachments = []
    for m in materials_list:
        if 'driveFile' in m:
            df = m['driveFile'].get('driveFile', {})
            mime = df.get('mimeType', '')
            if 'pdf' in mime:
                ftype = 'pdf'
            elif 'document' in mime:
                ftype = 'document'
            elif 'presentation' in mime:
                ftype = 'presentation'
            elif 'spreadsheet' in mime:
                ftype = 'spreadsheet'
            elif mime.startswith('image/'):
                ftype = 'image'
            else:
                ftype = 'file'
            attachments.append({
                "type": "drive",
                "file_type": ftype,
                "title": df.get('title', 'Drive File'),
                "url": df.get('alternateLink', ''),
                "file_id": df.get('id', ''),
                "thumbnail": df.get('thumbnailUrl', ''),
                "mime_type": mime,
            })
        elif 'youtubeVideo' in m:
            yt = m['youtubeVideo']
            vid = yt.get('id', '')
            attachments.append({
                "type": "video",
                "file_type": "youtube",
                "title": yt.get('title', 'YouTube Video'),
                "url": yt.get('alternateLink', ''),
                "video_id": vid,
                "thumbnail": f"https://img.youtube.com/vi/{vid}/mqdefault.jpg" if vid else '',
            })
        elif 'link' in m:
            l = m['link']
            attachments.append({
                "type": "link",
                "file_type": "web",
                "title": l.get('title', 'Link'),
                "url": l.get('url', ''),
            })
        elif 'form' in m:
            f = m['form']
            attachments.append({
                "type": "form",
                "file_type": "google_form",
                "title": f.get('title', 'Form'),
                "url": f.get('formUrl', ''),
            })
    return attachments


# Global lock for Google API calls to prevent thread/async conflicts and rate limits
from app.utils.locks import GLOBAL_API_LOCK

async def _get_creds(user: User):
    """Build Google credentials from stored tokens, refreshing if necessary."""
    if not user.google_access_token:
        raise HTTPException(status_code=400, detail="Google account not connected")
    
    from google.auth.transport.requests import Request
    creds = build_credentials(user.google_access_token, user.google_refresh_token, user.google_token_expiry)
    
    # Proactively refresh if expired
    if creds.expired or (creds.expiry and creds.expiry < datetime.now(timezone.utc)):
        if user.google_refresh_token:
            try:
                await asyncio.to_thread(creds.refresh, Request())
                # Update user tokens in DB after refresh
                from app.models.database import AsyncSessionLocal
                async with AsyncSessionLocal() as db:
                    db_user = await db.get(User, user.id)
                    if db_user:
                        db_user.google_access_token = creds.token
                        db_user.google_token_expiry = creds.expiry
                        await db.commit()
                return creds
            except Exception as e:
                logger.error(f"Token refresh failed: {e}")
                raise HTTPException(status_code=401, detail="Session expired, please re-login")
    return creds


# ── Google Classroom ──────────────────────────────────────────────────────
@router.post("/sync/classroom")
async def sync_classroom(
    background_tasks: BackgroundTasks,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Sync Google Classroom courses and materials."""
    creds = await _get_creds(current_user)
    background_tasks.add_task(_do_classroom_sync, current_user.id, current_user.email, creds)
    return {"ok": True, "message": "Classroom sync started in background"}


async def _do_classroom_sync(user_id: int, email: str, creds):
    from googleapiclient.discovery import build
    from app.models.database import AsyncSessionLocal
    try:
        service = build("classroom", "v1", credentials=creds)
        drive_service = build("drive", "v3", credentials=creds) # Build drive service for attachments
        courses_resp = service.courses().list(courseStates=["ACTIVE"]).execute()
        courses_data = courses_resp.get("courses", [])

        async with AsyncSessionLocal() as db:
            print(f"🚀 [Sync] Starting Deep Classroom sync for {email}")
            for course_data in courses_data:
                # Upsert course
                result = await db.execute(
                    select(Course).where(Course.user_id == user_id, Course.external_id == course_data["id"])
                )
                course = result.scalar_one_or_none()
                teacher = course_data.get("ownerId", "")
                if not course:
                    course = Course(
                        user_id=user_id,
                        external_id=course_data["id"],
                        name=course_data.get("name", ""),
                        section=course_data.get("section"),
                        teacher=teacher,
                        source="classroom",
                    )
                    db.add(course)
                    await db.commit()
                    await db.refresh(course)
                
                print(f"📘 [Sync] Checking: {course.name}")

                # ── Fetch ALL CourseWork (Assignments) ──
                page_token = None
                while True:
                    await asyncio.sleep(0.5) # Rate limit safety
                    cw_resp = service.courses().courseWork().list(
                        courseId=course_data["id"],
                        pageSize=50,
                        pageToken=page_token
                    ).execute()
                    
                    course_works = cw_resp.get("courseWork", [])
                    print(f"   - Found batch of {len(course_works)} coursework/assignments")
                    
                    for cw in course_works:
                        title = cw.get('title', 'Untitled Assignment')
                        desc = cw.get('description','')
                        doc = f"Course: {course.name}\nType: Assignment\nTitle: {title}\nDescription: {desc}\nDue: {cw.get('dueDate', {})}"
                        
                        # 1. Save to local DB first to get ID
                        cw_check = await db.execute(
                            select(Material).where(Material.user_id == user_id, Material.external_id == cw["id"])
                        )
                        existing_cw = cw_check.scalar_one_or_none()
                        
                        cw_link = cw.get('alternateLink')
                        cw_attachments = _extract_classroom_attachments(cw.get('materials', []))
                        
                        if existing_cw:
                            mat_for_index = existing_cw
                            # Update link/attachments if missing
                            if not mat_for_index.source_link: mat_for_index.source_link = cw_link
                            if not mat_for_index.attachments: mat_for_index.attachments = json.dumps(cw_attachments)
                            await db.commit()
                        else:
                            mat_for_index = Material(
                                user_id=user_id, course_id=course.id,
                                title=title, material_type="assignment",
                                source="classroom", external_id=cw["id"],
                                content_preview=desc[:500],
                                full_text=desc,
                                indexed_in_memsapien=False,
                                source_link=cw_link,
                                attachments=json.dumps(cw_attachments),
                            )
                            db.add(mat_for_index)
                            await db.commit()
                            await db.refresh(mat_for_index)

                        # 2. Index assignment metadata with material_id
                        if not mat_for_index.indexed_in_memsapien:
                            sm_id = await get_ms_client().add_document(
                                content=doc, 
                                metadata={
                                    "source": "classroom", 
                                    "type": "academic_material", 
                                    "course": course.name, 
                                    "user_id": email,
                                    "material_id": str(mat_for_index.id),
                                    "source_link": cw_link
                                }, 
                                user_email=email, 
                                title=title,
                                container_tag=f"course_{course.id}"
                            )
                            if sm_id:
                                mat_for_index.indexed_in_memsapien = True
                                await db.commit()
                        
                        # 2. Deep Index Attachments & Embedded Links
                        cw_processed_ids = set()
                        
                        # A. Embedded Links in description
                        found_ids = re.findall(r'https?://(?:docs\.google\.com/(?:document|spreadsheets|presentation|forms)/d/|drive\.google\.com/(?:file/d/|open\?id=))([a-zA-Z0-9_-]+)', desc)
                        for fid in found_ids:
                            if fid not in cw_processed_ids:
                                # Check for duplicate embedded link
                                att_check = await db.execute(
                                    select(Material).where(Material.user_id == user_id, Material.external_id == fid)
                                )
                                if att_check.scalar_one_or_none():
                                    print(f"     ✅ [Skip] Embedded link {fid} already synced")
                                    cw_processed_ids.add(fid)
                                    continue

                                # EARLY DEDUPLICATION
                                dup_check = await db.execute(select(Material).where(Material.external_id == fid, Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Embedded file {fid} already synced & indexed")
                                    continue

                                print(f"     🔗 [Deep Sync] Found embedded Drive link in assignment: {fid}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, fid)
                                if content or raw_bytes:
                                    # Local DB
                                    new_mat = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Linked File ({fid[:8]})", material_type="file",
                                        source="classroom", external_id=fid,
                                        parent_external_id=cw["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content,
                                        file_content=raw_bytes,
                                        indexed_in_memsapien=True, # Will commit below
                                        source_link=f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/0"
                                    )
                                    db.add(new_mat)
                                    await db.commit()
                                    await db.refresh(new_mat)
                                    
                                    # Update with real ID
                                    new_mat.source_link = f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/{new_mat.id}"
                                    await db.commit()

                                    # Push to SM with correct link
                                    att_doc = f"Course: {course.name}\nSource Assignment: {title}\nFile Content:\n{content}"
                                    await get_ms_client().add_document(
                                        att_doc, 
                                        metadata={
                                            "source": "classroom", 
                                            "type": "academic_material", 
                                            "course": course.name, 
                                            "user_id": email,
                                            "material_id": str(new_mat.id),
                                            "source_link": new_mat.source_link
                                        }, 
                                        user_email=email, 
                                        title=f"Linked File ({fid[:8]})",
                                        container_tag=f"course_{course.id}"
                                    )

                        # B. Formal attachments
                        for mat_ref in cw.get('materials', []):
                            if 'driveFile' in mat_ref:
                                df = mat_ref['driveFile']['driveFile']
                                if df['id'] in cw_processed_ids: continue

                                # DEDUPLICATION: skip if indexed, upsert if exists but not indexed
                                dup_check = await db.execute(select(Material).where(Material.external_id == df['id'], Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Formal attachment {df['id']} already synced & indexed")
                                    cw_processed_ids.add(df['id'])
                                    continue

                                print(f"     📎 [Attachment] Found in assignment: {df['title']}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, df['id'], df.get('mimeType', ''))
                                backend_url = os.getenv('BACKEND_URL', 'http://localhost:8002')

                                if existing:
                                    existing.full_text = content
                                    existing.content_preview = content[:500] if content else existing.content_preview
                                    existing.file_content = raw_bytes
                                    existing.mime_type = df.get('mimeType')
                                    await db.commit()
                                    mat_to_index = existing
                                else:
                                    mat_to_index = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Attachment: {df['title']}", material_type="file",
                                        source="classroom", external_id=df["id"],
                                        parent_external_id=cw["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content if content else None,
                                        file_content=raw_bytes,
                                        mime_type=df.get('mimeType'),
                                        indexed_in_memsapien=False,
                                        source_link=f"{backend_url}/profile/proxy/drive/0"
                                    )
                                    db.add(mat_to_index)
                                    await db.commit()
                                    await db.refresh(mat_to_index)
                                    mat_to_index.source_link = f"{backend_url}/profile/proxy/drive/{mat_to_index.id}"
                                    await db.commit()

                                sm_content = f"Course: {course.name}\nSource Assignment: {title}\nFile: {df['title']}\n\n{content}" if content else f"Course: {course.name}\nSource Assignment: {title}\nFile attachment: {df['title']}\nType: {df.get('mimeType', 'unknown')}"
                                await get_ms_client().add_document(
                                    sm_content,
                                    metadata={
                                        "source": "classroom",
                                        "type": "academic_material",
                                        "course": course.name,
                                        "user_id": email,
                                        "material_id": str(mat_to_index.id),
                                        "source_link": mat_to_index.source_link
                                    },
                                    user_email=email,
                                    title=f"Attachment: {df['title']}",
                                    container_tag=f"course_{course.id}"
                                )
                                mat_to_index.indexed_in_memsapien = True
                                await db.commit()
                                cw_processed_ids.add(df['id'])
                            elif 'link' in mat_ref:
                                link_data = mat_ref['link']
                                link_url = link_data.get('url', '')
                                link_title = link_data.get('title', link_url)
                                if not link_url or link_url in cw_processed_ids: continue
                                print(f"     🔗 [Link] Found in assignment: {link_title}")
                                # Index to SM only — links are not stored in local DB
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nAssignment: {title}\nLinked resource: {link_title}\nURL: {link_url}",
                                    metadata={"source": "classroom", "type": "academic_material", "course": course.name, "user_id": email, "source_link": link_url},
                                    user_email=email, title=f"Link: {link_title}",
                                    container_tag=f"course_{course.id}"
                                )
                                cw_processed_ids.add(link_url)
                            elif 'youtubeVideo' in mat_ref:
                                yt = mat_ref['youtubeVideo']
                                yt_url = yt.get('alternateLink', '')
                                yt_title = yt.get('title', 'YouTube Video')
                                if not yt_url or yt_url in cw_processed_ids: continue
                                print(f"     ▶️ [YouTube] Found in assignment: {yt_title}")
                                new_mat = Material(
                                    user_id=user_id, course_id=course.id,
                                    title=f"Video: {yt_title}", material_type="file",
                                    source="classroom", external_id=None,
                                    parent_external_id=cw["id"],
                                    content_preview=f"YouTube video: {yt_title}",
                                    source_link=yt_url,
                                    indexed_in_memsapien=False,
                                )
                                db.add(new_mat)
                                await db.commit()
                                await db.refresh(new_mat)
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nAssignment: {title}\nYouTube video: {yt_title}\nURL: {yt_url}",
                                    metadata={"source": "classroom", "type": "academic_material", "course": course.name, "user_id": email, "source_link": yt_url},
                                    user_email=email, title=f"Video: {yt_title}",
                                    container_tag=f"course_{course.id}"
                                )
                                new_mat.indexed_in_memsapien = True
                                await db.commit()
                                cw_processed_ids.add(yt_url)

                        cw_alternate_link = cw.get('alternateLink')
                        cw_attachments = _extract_classroom_attachments(cw.get('materials', []))
                        if existing_cw:
                            updated = False
                            if not existing_cw.source_link and cw_alternate_link:
                                existing_cw.source_link = cw_alternate_link
                                updated = True
                            if (not existing_cw.attachments or existing_cw.attachments == '[]') and cw_attachments:
                                existing_cw.attachments = json.dumps(cw_attachments)
                                updated = True
                            if updated:
                                await db.commit()
                            print(f"   ✅ [Skip] Assignment {cw['id']} already synced")
                        else:
                            db.add(Material(
                                user_id=user_id, course_id=course.id,
                                title=title, material_type="assignment",
                                source="classroom", external_id=cw["id"],
                                content_preview=desc[:500],
                                full_text=desc,
                                indexed_in_memsapien=bool(sm_id),
                                source_link=cw_alternate_link,
                                attachments=json.dumps(cw_attachments),
                            ))
                        # Pagination Check
                    page_token = cw_resp.get("nextPageToken")
                    if not page_token:
                        break

                # ── Fetch ALL Announcements ──
                page_token = None
                while True:
                    await asyncio.sleep(0.5)
                    ann_resp = service.courses().announcements().list(
                        courseId=course_data["id"],
                        pageSize=50,
                        pageToken=page_token
                    ).execute()
                    
                    announcements = ann_resp.get("announcements", [])
                    print(f"   - Found batch of {len(announcements)} announcements")
                    
                    for ann in announcements:
                        await asyncio.sleep(0.5)
                        text = ann.get('text','')
                        doc = f"Course: {course.name}\nType: Announcement\nText: {text}"
                        
                        # 1. Save locally first to get ID
                        ann_id_raw = ann["id"]
                        ann_check = await db.execute(
                            select(Material).where(Material.user_id == user_id, Material.external_id == ann_id_raw)
                        )
                        existing_ann = ann_check.scalar_one_or_none()
                        
                        ann_link = ann.get('alternateLink')
                        ann_attachments = _extract_classroom_attachments(ann.get('materials', []))
                        
                        if existing_ann:
                            ann_mat = existing_ann
                            # Backfill missing info
                            if not ann_mat.source_link: ann_mat.source_link = ann_link
                            if not ann_mat.attachments: ann_mat.attachments = json.dumps(ann_attachments)
                            await db.commit()
                        else:
                            ann_mat = Material(
                                user_id=user_id, course_id=course.id,
                                title=f"Announcement: {course.name}", material_type="announcement",
                                source="classroom", external_id=ann_id_raw,
                                content_preview=text[:500],
                                full_text=text,
                                indexed_in_memsapien=False,
                                source_link=ann_link,
                                attachments=json.dumps(ann_attachments),
                                created_at=datetime.fromisoformat(ann.get("updateTime", ann.get("creationTime")).replace("Z", "+00:00")) if ann.get("creationTime") else datetime.now(timezone.utc),
                            )
                            db.add(ann_mat)
                            await db.commit()
                            await db.refresh(ann_mat)

                        # 2. Index the announcement text with material_id
                        if not ann_mat.indexed_in_memsapien:
                            sm_id = await get_ms_client().add_document(
                                doc, 
                                metadata={
                                    "source": "classroom", 
                                    "type": "academic_material", 
                                    "course": course.name, 
                                    "user_id": email,
                                    "material_id": str(ann_mat.id),
                                    "source_link": ann_link
                                }, 
                                user_email=email, 
                                title=f"Announcement: {course.name}",
                                container_tag=f"course_{course.id}"
                            )
                            if sm_id:
                                ann_mat.indexed_in_memsapien = True
                                await db.commit()
                        # 2. Deep Index attachments & Embedded Links
                        processed_file_ids = set()
                        
                        # A. Formal materials
                        all_materials = ann.get('materials', [])
                        
                        # B. Embedded Links in text
                        found_ids = re.findall(r'https?://(?:docs\.google\.com/(?:document|spreadsheets|presentation|forms)/d/|drive\.google\.com/(?:file/d/|open\?id=))([a-zA-Z0-9_-]+)', text)
                        for fid in found_ids:
                            if fid not in processed_file_ids:
                                # EARLY DEDUPLICATION
                                dup_check = await db.execute(select(Material).where(Material.external_id == fid, Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Announcement link {fid} already synced & indexed")
                                    processed_file_ids.add(fid)
                                    continue

                                print(f"     🔗 [Deep Sync] Found embedded Drive link in text: {fid}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, fid)
                                if content or raw_bytes:
                                    # Save to local DB
                                    new_mat = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Linked File ({fid[:8]})", material_type="file",
                                        source="classroom", external_id=fid,
                                        parent_external_id=ann["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content,
                                        file_content=raw_bytes,
                                        indexed_in_memsapien=True, # Will commit below
                                        source_link=f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/0"
                                    )
                                    db.add(new_mat)
                                    await db.commit()
                                    await db.refresh(new_mat)

                                    # Update with real ID and link
                                    new_mat.source_link = f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/{new_mat.id}"
                                    await db.commit()

                                    # Push to SM
                                    att_doc = f"Course: {course.name}\nSource Announcement: {text[:200]}\nFile Content:\n{content}"
                                    await get_ms_client().add_document(
                                        att_doc, 
                                        metadata={
                                            "source": "classroom", 
                                            "type": "academic_material", 
                                            "course": course.name, 
                                            "user_id": email,
                                            "file_id": fid,
                                            "material_id": str(new_mat.id),
                                            "source_link": new_mat.source_link
                                        }, 
                                        user_email=email,
                                        title=f"Drive Link in Announcement",
                                        container_tag=f"course_{course.id}"
                                    )
                                    processed_file_ids.add(fid)

                        # C. Proper attachments
                        for mat_ref in all_materials:
                            if 'driveFile' in mat_ref:
                                df = mat_ref['driveFile']['driveFile']
                                if df['id'] in processed_file_ids: continue

                                # DEDUPLICATION: skip if indexed, upsert if exists but not indexed
                                dup_check = await db.execute(select(Material).where(Material.external_id == df['id'], Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Announcement attachment {df['id']} already synced & indexed")
                                    processed_file_ids.add(df['id'])
                                    continue

                                print(f"     📎 [Attachment] Found in announcement: {df['title']}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, df['id'], df.get('mimeType', ''))
                                backend_url = os.getenv('BACKEND_URL', 'http://localhost:8002')

                                if existing:
                                    existing.full_text = content
                                    existing.content_preview = content[:500] if content else existing.content_preview
                                    existing.file_content = raw_bytes
                                    existing.mime_type = df.get('mimeType')
                                    await db.commit()
                                    mat_to_index = existing
                                else:
                                    mat_to_index = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Attachment: {df['title']}", material_type="file",
                                        source="classroom", external_id=df["id"],
                                        parent_external_id=ann["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content if content else None,
                                        file_content=raw_bytes,
                                        mime_type=df.get('mimeType'),
                                        indexed_in_memsapien=False,
                                        source_link=f"{backend_url}/profile/proxy/drive/0"
                                    )
                                    db.add(mat_to_index)
                                    await db.commit()
                                    await db.refresh(mat_to_index)
                                    mat_to_index.source_link = f"{backend_url}/profile/proxy/drive/{mat_to_index.id}"
                                    await db.commit()

                                sm_content = f"Course: {course.name}\nSource Announcement: {text[:100]}\nFile: {df['title']}\n\n{content}" if content else f"Course: {course.name}\nSource Announcement: {text[:100]}\nFile attachment: {df['title']}\nType: {df.get('mimeType', 'unknown')}"
                                await get_ms_client().add_document(
                                    sm_content,
                                    metadata={
                                        "source": "classroom",
                                        "type": "academic_material",
                                        "course": course.name,
                                        "user_id": email,
                                        "file_id": df["id"],
                                        "material_id": str(mat_to_index.id),
                                        "source_link": mat_to_index.source_link
                                    },
                                    user_email=email,
                                    title=f"Attachment: {df['title']}",
                                    container_tag=f"user_{user_id}"
                                )
                                mat_to_index.indexed_in_memsapien = True
                                await db.commit()
                                processed_file_ids.add(df['id'])
                            elif 'link' in mat_ref:
                                link_data = mat_ref['link']
                                link_url = link_data.get('url', '')
                                link_title = link_data.get('title', link_url)
                                if not link_url or link_url in processed_file_ids: continue
                                print(f"     🔗 [Link] Found in announcement: {link_title}")
                                # Index to SM only — links are not stored in local DB
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nAnnouncement: {text[:100]}\nLinked resource: {link_title}\nURL: {link_url}",
                                    metadata={"source": "classroom", "type": "academic_material", "course": course.name, "user_id": email, "source_link": link_url},
                                    user_email=email, title=f"Link: {link_title}", container_tag=f"user_{user_id}"
                                )
                                processed_file_ids.add(link_url)
                            elif 'youtubeVideo' in mat_ref:
                                yt = mat_ref['youtubeVideo']
                                yt_url = yt.get('alternateLink', '')
                                yt_title = yt.get('title', 'YouTube Video')
                                if not yt_url or yt_url in processed_file_ids: continue
                                print(f"     ▶️ [YouTube] Found in announcement: {yt_title}")
                                new_mat = Material(
                                    user_id=user_id, course_id=course.id,
                                    title=f"Video: {yt_title}", material_type="file",
                                    source="classroom", external_id=None,
                                    parent_external_id=ann["id"],
                                    content_preview=f"YouTube video: {yt_title}",
                                    source_link=yt_url,
                                    indexed_in_memsapien=False,
                                )
                                db.add(new_mat)
                                await db.commit()
                                await db.refresh(new_mat)
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nAnnouncement: {text[:100]}\nYouTube video: {yt_title}\nURL: {yt_url}",
                                    metadata={"source": "classroom", "type": "academic_material", "course": course.name, "user_id": email, "source_link": yt_url},
                                    user_email=email, title=f"Video: {yt_title}", container_tag=f"user_{user_id}"
                                )
                                new_mat.indexed_in_memsapien = True
                                await db.commit()
                                processed_file_ids.add(yt_url)
                        
                    # Pagination Check
                    page_token = ann_resp.get("nextPageToken")
                    if not page_token:
                        break

                # ── Fetch ALL Materials (Resources) ──
                page_token = None
                while True:
                    await asyncio.sleep(0.5)
                    mat_resp = service.courses().courseWorkMaterials().list(
                        courseId=course_data["id"],
                        pageSize=50,
                        pageToken=page_token
                    ).execute()
                    
                    course_materials = mat_resp.get("courseWorkMaterial", [])
                    print(f"   - Found batch of {len(course_materials)} materials/resources")
                    
                    for cm in course_materials:
                        await asyncio.sleep(0.5)
                        title = cm.get('title', 'Untitled Resource')
                        desc = cm.get('description','')
                        doc = f"Course: {course.name}\nType: Resource\nTitle: {title}\nDescription: {desc}"
                        
                        # 1. Save locally first to get ID
                        cm_id_raw = cm["id"]
                        cm_check = await db.execute(
                            select(Material).where(Material.user_id == user_id, Material.external_id == cm_id_raw)
                        )
                        existing_cm = cm_check.scalar_one_or_none()
                        
                        cm_link = cm.get('alternateLink')
                        cm_attachments = _extract_classroom_attachments(cm.get('materials', []))
                        
                        if existing_cm:
                            cm_mat = existing_cm
                            # Backfill
                            if not cm_mat.source_link: cm_mat.source_link = cm_link
                            if not cm_mat.attachments: cm_mat.attachments = json.dumps(cm_attachments)
                            await db.commit()
                        else:
                            cm_mat = Material(
                                user_id=user_id, course_id=course.id,
                                title=title, material_type="material",
                                source="classroom", external_id=cm_id_raw,
                                content_preview=desc[:500],
                                full_text=desc,
                                indexed_in_memsapien=False,
                                source_link=cm_link,
                                attachments=json.dumps(cm_attachments),
                            )
                            db.add(cm_mat)
                            await db.commit()
                            await db.refresh(cm_mat)

                        # 2. Index the resource metadata with material_id
                        if not cm_mat.indexed_in_memsapien:
                            sm_id = await get_ms_client().add_document(
                                doc, 
                                metadata={
                                    "source": "classroom", 
                                    "type": "academic_material", 
                                    "course": course.name, 
                                    "user_id": email,
                                    "material_id": str(cm_mat.id),
                                    "source_link": cm_link
                                }, 
                                user_email=email, 
                                title=title
                            )
                            if sm_id:
                                cm_mat.indexed_in_memsapien = True
                                await db.commit()

                        # 3. Deep Index attachments & Embedded Links
                        resource_processed_ids = set()
                        
                        # A. Embedded Links in description
                        found_ids = re.findall(r'https?://(?:docs\.google\.com/(?:document|spreadsheets|presentation|forms)/d/|drive\.google\.com/(?:file/d/|open\?id=))([a-zA-Z0-9_-]+)', desc)
                        for fid in found_ids:
                            if fid not in resource_processed_ids:
                                # EARLY DEDUPLICATION
                                dup_check = await db.execute(select(Material).where(Material.external_id == fid, Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Resource link {fid} already synced & indexed")
                                    resource_processed_ids.add(fid)
                                    continue

                                print(f"     🔗 [Deep Sync] Found embedded Drive link in resource: {fid}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, fid)
                                if content or raw_bytes:
                                    # Save to local DB
                                    new_mat = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Linked File ({fid[:8]})", material_type="file",
                                        source="classroom", external_id=fid,
                                        parent_external_id=cm["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content,
                                        file_content=raw_bytes,
                                        indexed_in_memsapien=True, # Will commit below
                                        source_link=f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/0"
                                    )
                                    db.add(new_mat)
                                    await db.commit()
                                    await db.refresh(new_mat)

                                    # Update with real ID and link
                                    new_mat.source_link = f"{os.getenv('BACKEND_URL', 'http://localhost:8002')}/profile/proxy/drive/{new_mat.id}"
                                    await db.commit()

                                    # Push to SM
                                    att_doc = f"Course: {course.name}\nSource Resource: {title}\nFile Content:\n{content}"
                                    await get_ms_client().add_document(
                                        att_doc, 
                                        metadata={
                                            "source": "classroom", 
                                            "type": "academic_material", 
                                            "course": course.name, 
                                            "user_id": email,
                                            "file_id": fid,
                                            "material_id": str(new_mat.id),
                                            "source_link": new_mat.source_link
                                        }, 
                                        user_email=email, 
                                        title=f"Linked File ({fid[:8]})",
                                        container_tag=f"user_{user_id}"
                                    )
                                    resource_processed_ids.add(fid)

                        # B. Formal attachments
                        for mat_ref in cm.get('materials', []):
                            if 'driveFile' in mat_ref:
                                df = mat_ref['driveFile']['driveFile']
                                if df['id'] in resource_processed_ids: continue

                                # DEDUPLICATION: skip if indexed, upsert if exists but not indexed
                                dup_check = await db.execute(select(Material).where(Material.external_id == df['id'], Material.user_id == user_id))
                                existing = dup_check.scalars().first()
                                if existing and existing.indexed_in_memsapien:
                                    print(f"     ✅ [Skip] Resource attachment {df['id']} already synced & indexed")
                                    resource_processed_ids.add(df['id'])
                                    continue

                                print(f"     📎 [Attachment] Found in resource: {df['title']}")
                                content, raw_bytes = await _get_drive_file_content(drive_service, df['id'], df.get('mimeType', ''))
                                backend_url = os.getenv('BACKEND_URL', 'http://localhost:8002')

                                if existing:
                                    existing.full_text = content
                                    existing.content_preview = content[:500] if content else existing.content_preview
                                    existing.file_content = raw_bytes
                                    existing.mime_type = df.get('mimeType')
                                    await db.commit()
                                    mat_to_index = existing
                                else:
                                    mat_to_index = Material(
                                        user_id=user_id, course_id=course.id,
                                        title=f"Attachment: {df['title']}", material_type="file",
                                        source="classroom", external_id=df["id"],
                                        parent_external_id=cm["id"],
                                        content_preview=content[:500] if content else None,
                                        full_text=content if content else None,
                                        file_content=raw_bytes,
                                        mime_type=df.get('mimeType'),
                                        indexed_in_memsapien=False,
                                        source_link=f"{backend_url}/profile/proxy/drive/0"
                                    )
                                    db.add(mat_to_index)
                                    await db.commit()
                                    await db.refresh(mat_to_index)
                                    mat_to_index.source_link = f"{backend_url}/profile/proxy/drive/{mat_to_index.id}"
                                    await db.commit()

                                sm_content = f"Course: {course.name}\nSource Resource: {title}\nFile: {df['title']}\n\n{content}" if content else f"Course: {course.name}\nSource Resource: {title}\nFile attachment: {df['title']}\nType: {df.get('mimeType', 'unknown')}"
                                await get_ms_client().add_document(
                                    sm_content,
                                    metadata={
                                        "source": "classroom",
                                        "type": "academic_material",
                                        "course": course.name,
                                        "user_id": email,
                                        "file_id": df["id"],
                                        "material_id": str(mat_to_index.id),
                                        "source_link": mat_to_index.source_link
                                    },
                                    user_email=email,
                                    title=f"Attachment: {df['title']}",
                                    container_tag=f"user_{user_id}"
                                )
                                mat_to_index.indexed_in_memsapien = True
                                await db.commit()
                                resource_processed_ids.add(df['id'])
                            elif 'link' in mat_ref:
                                link_data = mat_ref['link']
                                link_url = link_data.get('url', '')
                                link_title = link_data.get('title', link_url)
                                if not link_url or link_url in resource_processed_ids: continue
                                print(f"     🔗 [Link] Found in resource: {link_title}")
                                # Index to SM only — links are not stored in local DB
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nResource: {title}\nLinked resource: {link_title}\nURL: {link_url}",
                                    metadata={"source": "classroom", "type": "academic_material", "course": course.name, "user_id": email, "source_link": link_url},
                                    user_email=email, title=f"Link: {link_title}", container_tag=f"user_{user_id}"
                                )
                                resource_processed_ids.add(link_url)
                            elif 'youtubeVideo' in mat_ref:
                                yt = mat_ref['youtubeVideo']
                                yt_url = yt.get('alternateLink', '')
                                yt_title = yt.get('title', 'YouTube Video')
                                if not yt_url or yt_url in resource_processed_ids: continue
                                print(f"     ▶️ [YouTube] Found in resource: {yt_title}")
                                new_mat = Material(
                                    user_id=user_id, course_id=course.id,
                                    title=f"Video: {yt_title}", material_type="file",
                                    source="classroom", external_id=None,
                                    parent_external_id=cm["id"],
                                    content_preview=f"YouTube video: {yt_title}",
                                    source_link=yt_url,
                                    indexed_in_memsapien=False,
                                )
                                db.add(new_mat)
                                await db.commit()
                                await get_ms_client().add_document(
                                    f"Course: {course.name}\nResource: {title}\nYouTube video: {yt_title}\nURL: {yt_url}",
                                    metadata={
                                        "source": "classroom", 
                                        "type": "academic_material", 
                                        "course": course.name, 
                                        "user_id": email, 
                                        "material_id": str(new_mat.id),
                                        "source_link": yt_url
                                    },
                                    user_email=email, title=f"Video: {yt_title}", container_tag=f"user_{user_id}"
                                )
                                new_mat.indexed_in_memsapien = True
                                await db.commit()
                                resource_processed_ids.add(yt_url)
                        
                        await db.commit() # Save incrementally

            # Mark user as synced
            result = await db.execute(select(User).where(User.id == user_id))
            user = result.scalar_one_or_none()
            if user:
                user.classroom_synced = True
            await db.commit()
    except Exception as e:
        print(f"[Classroom Sync] Error: {e}")




# ── Google Drive ──────────────────────────────────────────────────────────
async def _get_drive_file_content(service, file_id: str, mime_type: str = "") -> tuple[str, Optional[bytes]]:
    """Helper to extract text and binary from Google Drive files for Classroom sync."""
    async with GLOBAL_API_LOCK:
        try:
            if not mime_type:
                meta = await asyncio.to_thread(lambda: service.files().get(fileId=file_id, fields="mimeType").execute())
                mime_type = meta.get("mimeType", "")

            raw_bytes = await asyncio.to_thread(lambda: service.files().get_media(fileId=file_id).execute())
            
            if mime_type == "application/vnd.google-apps.document":
                export = await asyncio.to_thread(lambda: service.files().export(fileId=file_id, mimeType="text/plain").execute())
                return export.decode("utf-8", errors="replace")[:15000], raw_bytes
            elif mime_type == "application/pdf":
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(raw_bytes))
                text = "\n".join(p.extract_text() or "" for p in reader.pages)[:15000]
                return text, raw_bytes
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                import docx
                doc = docx.Document(io.BytesIO(raw_bytes))
                text = "\n".join(para.text for para in doc.paragraphs)[:15000]
                return text, raw_bytes
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw_bytes))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return "\n".join(text_runs)[:15000], raw_bytes
            elif mime_type.startswith("image/"):
                # Vision OCR using gpt-4o-mini
                try:
                    b64 = base64.b64encode(raw_bytes).decode("utf-8")
                    from app.services.openai_service import unified_completion
                    resp = await unified_completion(
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": "Extract all text from this image and provide a high-level academic description of what it shows. Return the transcription followed by the description."},
                                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}}
                                ]
                            }
                        ],
                        tier="flash"
                    )
                    return resp["text"][:15000], raw_bytes
                except Exception as ve:
                    logger.error(f"Vision OCR error: {ve}")
                    return "[Image content could not be processed]", raw_bytes
            elif "google-apps" in mime_type:
                try:
                    export = await asyncio.to_thread(lambda: service.files().export(fileId=file_id, mimeType="text/plain").execute())
                    return export.decode("utf-8", errors="replace")[:10000], raw_bytes
                except: return "", raw_bytes
            elif mime_type == "text/plain":
                return raw_bytes.decode("utf-8", errors="replace")[:10000], raw_bytes
            else:
                return "", raw_bytes
        except Exception as e:
            logger.error(f"Classroom file extraction error: {e}")
            return "", None




@router.get("/drive/proxy/{file_id}")
async def drive_proxy_legacy(file_id: str, db: AsyncSession = Depends(get_db)):
    """Legacy compatibility: redirect old raw File IDs to new localized Material IDs."""
    result = await db.execute(
        select(Material).where(Material.external_id == file_id).limit(1)
    )
    mat = result.scalar_one_or_none()
    if mat:
        backend_url = os.getenv("BACKEND_URL", "http://localhost:8002")
        return RedirectResponse(url=f"{backend_url}/profile/proxy/drive/{mat.id}")
    
    raise HTTPException(status_code=404, detail="Source material not found or not synced.")


@router.post("/drive/sync")
async def sync_drive(
    request: DriveSyncRequest,
    background_tasks: BackgroundTasks,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Sync specifically selected Drive files in the background."""
    service = DriveSyncService(current_user)
    
    if request.file_ids:
        background_tasks.add_task(service.sync_selected_files, request.file_ids)
        return {"message": "File sync started", "count": len(request.file_ids)}
    
    return {"error": "No files specified"}

@router.get("/drive/config")
async def get_drive_config(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """
    Returns API keys and the current user's email for the Google Picker.
    Also returns the existing access token if valid, to assist auto-login.
    """
    try:
        client_id = os.getenv("GOOGLE_CLIENT_ID", "")
        
        # Proactively refresh token if possible to ensure the picker has a valid one
        token = current_user.google_access_token
        if current_user.google_refresh_token:
            try:
                creds = await asyncio.to_thread(
                    refresh_credentials,
                    current_user.google_refresh_token
                )
                if creds and creds.token:
                    token = creds.token
                    current_user.google_access_token = token
                    await db.commit()
            except Exception as e:
                logger.error(f"Failed to refresh token for picker: {e}")

        return {
            "apiKey": os.getenv("GOOGLE_PICKER_API_KEY"),
            "clientId": client_id,
            "appId": client_id.split("-")[0] if client_id else "",
            "userEmail": current_user.email,
            "accessToken": token
        }
    except Exception as e:
        logger.error(f"Failed to produce drive config: {e}")
        return {"error": str(e)}

# Note: legacy /drive/proxy/{file_id} removed in favor of localized /profile/proxy/drive/{material_id}


# ── Status ────────────────────────────────────────────────────────────────
@router.get("/status")
async def integration_status(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    courses_res = await db.execute(select(Course).where(Course.user_id == current_user.id))
    courses = courses_res.scalars().all()
    materials_res = await db.execute(select(Material).where(Material.user_id == current_user.id))
    materials = materials_res.scalars().all()
    return {
        "classroom": {"synced": current_user.classroom_synced, "courses": len([c for c in courses if c.source == "classroom"])},
        "drive": {"synced": current_user.drive_synced, "files": len([m for m in materials if m.source == "drive"])},
        "total_materials": len(materials),
    }
