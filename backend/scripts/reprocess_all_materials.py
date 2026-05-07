import asyncio
import io
import base64
import os
import json
import logging
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession, async_sessionmaker
from sqlalchemy import select
from app.models.database import Material, User
from docx import Document
from pptx import Presentation
from app.services.openai_service import get_client

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

DATABASE_URL = "sqlite+aiosqlite:///backend/fiwb.db"

async def extract_content(mime_type, raw_bytes):
    if not raw_bytes:
        return ""
    
    try:
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(io.BytesIO(raw_bytes))
            return "\n".join(para.text for para in doc.paragraphs)
        
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(raw_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        
        elif mime_type.startswith("image/"):
            b64 = base64.b64encode(raw_bytes).decode("utf-8")
            client = get_client()
            resp = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract all text from this image and provide a high-level academic description of what it shows. Return the transcription followed by the description."},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}}
                        ]
                    }
                ],
                max_tokens=1000,
            )
            return resp.choices[0].message.content
    except Exception as e:
        logger.error(f"Extraction error for {mime_type}: {e}")
        return ""
    return ""

async def reprocess_all():
    engine = create_async_engine(DATABASE_URL)
    async_session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)
    
    async with async_session() as session:
        # Find all materials with empty full_text
        stmt = select(Material).where(
            (Material.full_text == None) | (Material.full_text == "")
        )
        result = await session.execute(stmt)
        materials = result.scalars().all()
        
        logger.info(f"Found {len(materials)} materials with empty text.")
        
        for mat in materials:
            mtype = mat.mime_type or ""
            title = mat.title.lower()
            content = mat.file_content
            
            # Smart detection for mislabeled or missing MIME types
            if content and content.startswith(b"PK\x03\x04"):
                if title.endswith(".docx"): mtype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                elif title.endswith(".pptx"): mtype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            elif content and content.startswith(b"%PDF"):
                mtype = "application/pdf"
            elif title.endswith(".jpg") or title.endswith(".jpeg"): mtype = "image/jpeg"
            elif title.endswith(".png"): mtype = "image/png"
            
            if not mtype or mtype == "application/octet-stream":
                # Fallback to extension check if still unknown
                if title.endswith(".docx"): mtype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                elif title.endswith(".pptx"): mtype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                elif title.endswith(".pdf"): mtype = "application/pdf"

            if mtype in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/pdf"
            ] or mtype.startswith("image/"):
                logger.info(f"Reprocessing Material {mat.id}: {mat.title} (Detected: {mtype})")
                
                # Special handling for PDF since the main script handles it but it might have been missed
                if mtype == "application/pdf":
                    from pypdf import PdfReader
                    try:
                        reader = PdfReader(io.BytesIO(content))
                        text = "\n".join(p.extract_text() or "" for p in reader.pages)
                    except: text = ""
                else:
                    text = await extract_content(mtype, content)
                
                if text:
                    mat.full_text = text
                    mat.content_preview = text[:500]
                    mat.mime_type = mtype # Correct it in DB too
                    logger.info(f"  Success! Extracted {len(text)} characters.")
                else:
                    logger.warning(f"  Failed to extract text for material {mat.id}.")
        
        await session.commit()
        logger.info("Done.")

if __name__ == "__main__":
    # Ensure OPENAI_API_KEY is available (from .env)
    from dotenv import load_dotenv
    load_dotenv("backend/.env")
    
    asyncio.run(reprocess_all())
